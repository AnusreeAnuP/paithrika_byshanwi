import os
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt

def generate_report():
    print("Generating Word Document...")
    doc = Document()
    
    # helper for headings
    def add_heading(text, level=1):
        h = doc.add_heading(text, level=level)
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER if level == 0 else WD_ALIGN_PARAGRAPH.LEFT

    # 1. Cover Page
    add_heading("PAITHRIKA_BYSHANWI", 0)
    doc.add_paragraph("\n\n\n\n\n")
    p = doc.add_paragraph("A Final Year Project Report")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph("\n\n")
    p = doc.add_paragraph("Submitted by:\n[Your Name]\n[Roll Number]")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()

    # 2. Certificate
    add_heading("CERTIFICATE")
    doc.add_paragraph("This is to certify that the project entitled 'Paithrika_ByShanwi' is a bonafide work carried out by [Your Name] in partial fulfillment of the requirements for the degree of Bachelor of Technology...")
    doc.add_page_break()

    # 3. Acknowledgement
    add_heading("ACKNOWLEDGEMENT")
    doc.add_paragraph("I would like to express my deep sense of gratitude to my project guide [Guide Name] and the Head of Department...")
    doc.add_page_break()

    # 4. Abstract
    add_heading("ABSTRACT")
    doc.add_paragraph("Paithrika_ByShanwi is a robust e-commerce platform developed using Python Django and PostgreSQL. The application focuses on the niche market of traditional ethnic wear (Sarees, Kurtis, etc.). It features an intuitive customer interface, a secure checkout process, and a comprehensive admin dashboard for sales management.")
    doc.add_page_break()

    # 5. Introduction
    add_heading("1. INTRODUCTION")
    doc.add_paragraph("This project aims to bridge the gap between traditional artisans and modern consumers through a specialized e-commerce experience...")
    
    add_heading("1.1 Problem Statement", 2)
    doc.add_paragraph("The existing manual systems for managing ethnic wear sales are inefficient, lacking scalability and global reach...")

    add_heading("1.2 Objectives", 2)
    doc.add_paragraph("- To provide a seamless online shopping experience for ethnic wear.\n- To automate inventory and order management.\n- To provide detailed sales analytics.")

    # 6. System Requirements
    add_heading("2. SYSTEM REQUIREMENTS")
    add_heading("2.1 Hardware Requirements", 2)
    doc.add_paragraph("- Processor: Intel i5 or higher\n- RAM: 8GB minimum\n- Storage: 256GB SSD")
    add_heading("2.2 Software Requirements", 2)
    doc.add_paragraph("- OS: Windows 10/11\n- Language: Python 3.13+\n- Framework: Django 5.0\n- Database: PostgreSQL\n- Frontend: HTML5, CSS3, JavaScript, Bootstrap 5")

    # 7. System Design
    add_heading("3. SYSTEM DESIGN")
    doc.add_paragraph("The system follows a Model-View-Template (MVT) architecture...")
    add_heading("3.1 ER Diagram", 2)
    doc.add_paragraph("[Insert ER Diagram Here - Describing Users, Products, Categories, Carts, and Orders]")
    
    # 8. Database Design
    add_heading("4. DATABASE DESIGN")
    doc.add_paragraph("The database consists of the following primary tables: CustomUser, Product, Category, Cart, CartItem, Order, and OrderItem.")

    # 9. Implementation & Testing
    add_heading("5. IMPLEMENTATION & TESTING")
    doc.add_paragraph("The codebase is modularized into several Django apps: accounts, products, cart, and orders...")
    add_heading("5.1 Test Cases", 2)
    table = doc.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Test ID'
    hdr_cells[1].text = 'Description'
    hdr_cells[2].text = 'Result'
    # sample row
    row_cells = table.add_row().cells
    row_cells[0].text = 'TC01'
    row_cells[1].text = 'User Registration'
    row_cells[2].text = 'Pass'

    # 10. Conclusion
    add_heading("6. CONCLUSION")
    doc.add_paragraph("Paithrika_ByShanwi successfully demonstrates the integration of modern web technologies to solve real-world retail challenges in the ethnic wear sector.")

    doc.save("Project_Report_Paithrika.docx")
    print("Word Document saved as 'Project_Report_Paithrika.docx'")

def generate_pptx():
    print("Generating PowerPoint Presentation...")
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Paithrika_ByShanwi"
    subtitle.text = "Ethnic Dress E-Commerce Platform\nFinal Year Project"

    # Slide 2: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technology Stack"
    content = slide.placeholders[1]
    content.text = "- Backend: Python Django\n- Frontend: HTML/CSS/JS/Bootstrap\n- Database: PostgreSQL\n- Tools: VS Code, Git"

    # Slide 3: Modules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Modules"
    content = slide.placeholders[1]
    content.text = "- User Authentication\n- Product Management\n- Shopping Cart & Wishlist\n- Order Checkout (COD)\n- Admin Sales Dashboard"

    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Features"
    content = slide.placeholders[1]
    content.text = "- Responsive UI Architecture\n- Real-time Product Search\n- Dynamic Category Filtering\n- Automated Inventory Alerts"

    # Slide 5: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Scope"
    content = slide.placeholders[1]
    content.text = "- Payment Gateway Integration (Stripe/Razorpay)\n- AI-based Product Recommendations\n- Mobile Application Development"

    prs.save("Project_Presentation_Paithrika.pptx")
    print("PowerPoint saved as 'Project_Presentation_Paithrika.pptx'")

if __name__ == "__main__":
    generate_report()
    generate_pptx()
